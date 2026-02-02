import os
import sys
from io import BytesIO

# Optional imports with error handling
try:
    from pptx import Presentation
except ImportError: Presentation = None

try:
    from docx import Document
except ImportError: Document = None

try:
    import openpyxl
except ImportError: openpyxl = None

try:
    from pypdf import PdfReader
except ImportError: PdfReader = None

def extract_text_from_file(file_path: str) -> str:
    """Detects file type and extracts text accordingly."""
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == '.pptx':
            return _read_pptx(file_path)
        elif ext == '.docx':
            return _read_docx(file_path)
        elif ext == '.xlsx':
            return _read_xlsx(file_path)
        elif ext == '.pdf':
            return _read_pdf(file_path)
        elif ext in ['.txt', '.md', '.csv', '.json', '.xml', '.html', '.py', '.c', '.cpp', '.h', '.log', '.sh']:
            return _read_text(file_path)
        else:
            return f"Error: Unsupported file extension '{ext}'"
    except Exception as e:
        return f"Error reading file {file_path}: {str(e)}"

def _read_pptx(path):
    if not Presentation: return "Error: python-pptx library missing"
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def _read_docx(path):
    if not Document: return "Error: python-docx library missing"
    doc = Document(path)
    return "\n".join([para.text for para in doc.paragraphs])

def _read_xlsx(path):
    if not openpyxl: return "Error: openpyxl library missing"
    wb = openpyxl.load_workbook(path, data_only=True)
    text = []
    for sheet in wb.sheetnames:
        text.append(f"--- Sheet: {sheet} ---")
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            # Convert row to string, filtering None
            row_text = [str(cell) for cell in row if cell is not None]
            if row_text:
                text.append("\t".join(row_text))
    return "\n".join(text)

def _read_pdf(path):
    if not PdfReader: return "Error: pypdf library missing"
    reader = PdfReader(path)
    text = []
    for page in reader.pages:
        text.append(page.extract_text() or "")
    return "\n".join(text)

def _read_text(path):
    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()
